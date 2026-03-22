"""
Generate KBCS Mid-Semester Presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import json

# ── Colour palette ──────────────────────────────────────────
BLUE   = RGBColor(19, 62, 110)
TEAL   = RGBColor(30, 126, 126)
GREEN  = RGBColor(52, 143, 80)
GRAY   = RGBColor(74, 74, 74)
LIGHT  = RGBColor(235, 243, 249)
WHITE  = RGBColor(255, 255, 255)
BLACK  = RGBColor(0, 0, 0)
RED_C  = RGBColor(192, 57, 43)
ORANGE = RGBColor(230, 126, 34)
DARK_BG = RGBColor(19, 62, 110)
ACCENT  = RGBColor(30, 126, 126)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ────────────────────────────────────────

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_dark_bg(slide):
    add_bg(slide, DARK_BG)

def add_accent_bar(slide, top=Inches(0), height=Inches(0.08)):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, prs.slide_width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = TEAL
    shp.line.fill.background()

def add_bottom_bar(slide, text="KBCS: Karma-Based Credit Scheduler  |  Mid-Semester Presentation"):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.05), prs.slide_width, Inches(0.45))
    bar.fill.solid()
    bar.fill.fore_color.rgb = BLUE
    bar.line.fill.background()
    tf = bar.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

def add_slide_number(slide, num, total=22):
    txBox = slide.shapes.add_textbox(Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.35))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = f"{num}/{total}"
    p.font.size = Pt(9)
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.RIGHT

def add_title_box(slide, title, subtitle=None, left=Inches(0.7), top=Inches(0.3), width=Inches(11.5)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.7))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = BLUE
    # underline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + Inches(0.65), Inches(3.5), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()
    if subtitle:
        sub = slide.shapes.add_textbox(left, top + Inches(0.8), width, Inches(0.4))
        stf = sub.text_frame
        sp = stf.paragraphs[0]
        sp.text = subtitle
        sp.font.size = Pt(16)
        sp.font.color.rgb = GRAY
        sp.font.italic = True

def add_content_box(slide, bullets, left=Inches(0.7), top=Inches(1.6), width=Inches(5.5), height=Inches(4.8), font_size=Pt(16), spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, (text, level) in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.level = level
        p.font.size = font_size if level == 0 else Pt(font_size.pt - 2)
        p.font.color.rgb = GRAY if level == 0 else RGBColor(100, 100, 100)
        p.space_after = spacing
        if level == 0:
            p.font.bold = True

def add_info_card(slide, title, body, left, top, width=Inches(3.5), height=Inches(2.8), title_color=TEAL):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(245, 248, 252)
    card.line.color.rgb = RGBColor(200, 210, 225)
    card.line.width = Pt(1)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.space_after = Pt(8)
    for line in body:
        p2 = tf.add_paragraph()
        p2.text = line
        p2.font.size = Pt(13)
        p2.font.color.rgb = GRAY
        p2.space_after = Pt(4)

def make_table(slide, rows, left, top, col_widths, row_height=Inches(0.4)):
    n_rows = len(rows)
    n_cols = len(rows[0])
    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, sum(col_widths), row_height * n_rows)
    table = table_shape.table
    for c, w in enumerate(col_widths):
        table.columns[c].width = w
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = WHITE if r == 0 else GRAY
                paragraph.alignment = PP_ALIGN.CENTER
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 252) if r % 2 == 0 else WHITE
    return table


TOTAL_SLIDES = 22

# ════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_dark_bg(sl)
# accent line
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.15), prs.slide_width, Inches(0.06))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()

# Title
t = sl.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(11), Inches(1.5))
tf = t.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Karma-Based Credit Scheduler (KBCS)"
p.font.size = Pt(42); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "A P4-Based Fairness Mechanism for Heterogeneous TCP Flows"
p2.font.size = Pt(22); p2.font.color.rgb = RGBColor(150, 190, 230); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)

# Team info
info = sl.shapes.add_textbox(Inches(2.5), Inches(3.7), Inches(8.5), Inches(2.5))
itf = info.text_frame; itf.word_wrap = True
lines = [
    ("Mid-Semester Project Presentation", Pt(18), RGBColor(180, 200, 220)),
    ("", Pt(10), WHITE),
    ("Gaurav Jaiswal (IIT2023157)  ·  Ishan Chadha (IIT2023158)  ·  Sameer Prasad (IIT2023134)", Pt(16), WHITE),
    ("", Pt(8), WHITE),
    ("Indian Institute of Information Technology, Allahabad", Pt(15), RGBColor(150, 190, 230)),
    ("March 2026", Pt(14), RGBColor(120, 160, 200)),
]
for i, (txt, sz, clr) in enumerate(lines):
    pp = itf.paragraphs[0] if i == 0 else itf.add_paragraph()
    pp.text = txt; pp.font.size = sz; pp.font.color.rgb = clr; pp.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════
# SLIDE 2 — AGENDA
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 2, TOTAL_SLIDES)
add_title_box(sl, "Presentation Outline")

agenda_items = [
    "The Fairness Problem in Modern Networks",
    "TCP CCA Families & Buffer Competition",
    "Programmable Data Planes with P4",
    "Baseline System: P4air — Architecture & Results",
    "Our Proposal: KBCS — Core Idea & Algorithm",
    "KBCS Architecture & Processing Pipeline",
    "Implementation Status & Design Gaps",
    "Experimental Evaluation & Metrics",
    "Roadmap to End-Semester",
    "Key Takeaways",
]
for i, item in enumerate(agenda_items):
    y = Inches(1.5) + Inches(0.48) * i
    # number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.0), y, Inches(0.38), Inches(0.38))
    circ.fill.solid(); circ.fill.fore_color.rgb = TEAL; circ.line.fill.background()
    ctf = circ.text_frame
    cp = ctf.paragraphs[0]; cp.text = str(i+1); cp.font.size = Pt(13); cp.font.color.rgb = WHITE; cp.font.bold = True; cp.alignment = PP_ALIGN.CENTER
    ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # text
    tb = sl.shapes.add_textbox(Inches(1.6), y, Inches(9), Inches(0.38))
    tp = tb.text_frame.paragraphs[0]; tp.text = item; tp.font.size = Pt(17); tp.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════
# SLIDE 3 — THE FAIRNESS PROBLEM
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 3, TOTAL_SLIDES)
add_title_box(sl, "The Fairness Problem", "Why do some TCP flows starve while others dominate?")

add_content_box(sl, [
    ("Real networks carry diverse TCP algorithms simultaneously", 0),
    ("CUBIC, BBR, Vegas, Illinois — all competing for the same buffer", 1),
    ("Aggressive senders fill buffers, delay-sensitive senders back off", 0),
    ("Result: unfair bandwidth distribution at bottleneck links", 1),
    ("Traditional AQM (RED, CoDel) treats all flows identically", 0),
    ("Cannot distinguish the aggressor from the victim", 1),
    ("Key Question:", 0),
    ("Can a programmable switch detect & correct this unfairness", 1),
    ("in real-time, without modifying the end hosts?", 1),
], top=Inches(1.8), width=Inches(6.5), font_size=Pt(17))

# Right side — visual metaphor card
add_info_card(sl, "🏟  The Buffer Arena", [
    "CUBIC:  \"I'll fill the buffer until",
    "            packets drop — that's how",
    "            I find the link capacity.\"",
    "",
    "BBR:     \"I'll estimate BDP and stay",
    "            below the queue limit.\"",
    "",
    "Vegas:  \"Any RTT increase means",
    "            congestion — I'll back off.\"",
    "",
    "→ Same buffer, conflicting strategies",
    "→ CUBIC dominates every time",
], left=Inches(7.8), top=Inches(1.8), width=Inches(4.8), height=Inches(4.5))


# ════════════════════════════════════════════════════════════
# SLIDE 4 — CCA TAXONOMY
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 4, TOTAL_SLIDES)
add_title_box(sl, "TCP Congestion Control Families", "Four families with fundamentally different strategies")

rows = [
    ["CCA Family", "Signal Used", "Examples", "Aggressiveness"],
    ["Purely Loss-Based", "Packet loss", "CUBIC, Reno, BIC", "Highest — fills buffer until loss"],
    ["Loss-Delay Hybrid", "Loss + RTT", "Illinois, YeAH, Veno", "Medium — delay as secondary signal"],
    ["Delay-Based", "RTT increase", "Vegas, LoLa", "Lowest — backs off at first sign"],
    ["Model-Based", "BDP estimation", "BBR", "Variable — periodic probing"],
]
make_table(sl, rows, Inches(0.8), Inches(1.8), [Inches(2.5), Inches(2.5), Inches(3.5), Inches(4.3)], Inches(0.6))

# Key insight box
add_info_card(sl, "⚡ Core Conflict", [
    "When loss-based and delay-based flows",
    "share a FIFO queue, loss-based ALWAYS wins.",
    "",
    "Loss-based needs losses to slow down.",
    "Delay-based backs off before losses occur.",
    "→ Fundamentally asymmetric competition.",
], left=Inches(1.5), top=Inches(4.5), width=Inches(10), height=Inches(2.0), title_color=RED_C)


# ════════════════════════════════════════════════════════════
# SLIDE 5 — BUFFERBLOAT & WHY AQM FAILS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 5, TOTAL_SLIDES)
add_title_box(sl, "Bufferbloat & Why Traditional AQM Fails")

add_content_box(sl, [
    ("The Bufferbloat Cascade", 0),
    ("CUBIC fills the shared buffer → queuing delay rises for ALL", 1),
    ("BBR sees inflated RTT → reduces pacing rate", 1),
    ("Vegas detects delay → backs off entirely → starvation", 1),
    ("", 0),
    ("Why RED/CoDel Can't Help", 0),
    ("Flow-agnostic: same drop probability for all packets", 1),
    ("Drops a Vegas packet = punishes the victim", 1),
    ("Drops a CUBIC packet = exactly what CUBIC expects", 1),
    ("", 0),
    ("What We Need Instead", 0),
    ("Behavior-aware differentiation at the switch", 1),
    ("Protect well-behaved flows, penalize aggressors", 1),
], top=Inches(1.5), width=Inches(6.5), font_size=Pt(16))

add_info_card(sl, "📊 Literature Evidence", [
    "Without CCA awareness:",
    "  Jain's FI as low as 0.18 – 0.49",
    "  (P4air paper, Turkovic & Kuipers 2021)",
    "",
    "With CCA-aware scheduling:",
    "  Jain's FI improves to 0.87 – 0.96",
    "",
    "→ Strong motivation for in-switch fairness",
], left=Inches(7.8), top=Inches(1.5), width=Inches(4.8), height=Inches(4.0))


# ════════════════════════════════════════════════════════════
# SLIDE 6 — PROGRAMMABLE DATA PLANES & P4
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 6, TOTAL_SLIDES)
add_title_box(sl, "P4 & Programmable Data Planes", "Moving fairness logic INTO the switch")

add_content_box(sl, [
    ("P4: Protocol-Independent Packet Processing", 0),
    ("Switch pipeline defined by the operator, not the vendor", 1),
    ("Parser → Ingress → Traffic Mgr → Egress → Deparser", 1),
    ("", 0),
    ("Key P4 Capabilities We Leverage", 0),
    ("Registers: persistent per-flow state across packets", 1),
    ("Hash functions: CRC16 for flow identification", 1),
    ("Bit-shift ALU: efficient decay without floating point", 1),
    ("Priority queues: differentiated scheduling in BMv2", 1),
    ("", 0),
    ("Our Environment", 0),
    ("BMv2 simple_switch  +  Mininet  +  P4_16", 1),
    ("Dumbbell topology: 10 Mbps bottleneck, 5ms delay", 1),
], top=Inches(1.5), width=Inches(7.0), font_size=Pt(16))

# v1model pipeline visual
add_info_card(sl, "v1model Pipeline", [
    " ┌─────────────────┐",
    " │     Parser       │",
    " └────────┬─────────┘",
    "          ↓",
    " ┌─────────────────┐",
    " │ Ingress Control  │ ← KBCS logic here",
    " └────────┬─────────┘",
    "          ↓",
    " ┌─────────────────┐",
    " │ Traffic Manager  │ ← Priority queues",
    " └────────┬─────────┘",
    "          ↓",
    " ┌─────────────────┐",
    " │ Egress Control   │",
    " └─────────────────┘",
], left=Inches(8.5), top=Inches(1.5), width=Inches(4.2), height=Inches(5.0))


# ════════════════════════════════════════════════════════════
# SLIDE 7 — LITERATURE OVERVIEW
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 7, TOTAL_SLIDES)
add_title_box(sl, "Related Work at a Glance", "10 papers spanning 4 research areas")

rows = [
    ["Research Area", "Key Papers", "What They Contribute"],
    ["CCA-Aware Fairness", "P4air, P4CCI, Real-Time CCA ID", "In-switch flow classification & treatment"],
    ["Buffer/Queue Mgmt", "CCA-aware QM, P4 Buffer Mgmt", "Programmable queue policies for fairness"],
    ["Fair Queueing", "PFQ (Proactive Fair Queueing)", "Fairness + high utilization simultaneously"],
    ["Network Telemetry", "HINT (In-Band Telemetry)", "P4-driven congestion feedback via INT"],
    ["Multi-Agent CC", "MARL-CC", "RL-based cooperative fairness learning"],
    ["P4 Frameworks", "P5, P4-NEON", "Event-driven & SDN-based P4 policies"],
]
make_table(sl, rows, Inches(0.6), Inches(1.7), [Inches(2.8), Inches(4.2), Inches(5.2)], Inches(0.55))

add_info_card(sl, "Research Gap We Target", [
    "FIFO / RED / CoDel → too simple, flow-agnostic",
    "P4air / P4CCI → effective but complex",
    "  (13+ registers, RTT estimation, recirculation)",
    "",
    "KBCS targets the middle ground:",
    "  Behavior-aware  +  Implementation-simple",
    "  → Only 2 registers per flow",
    "  → No timer dependency",
], left=Inches(1.5), top=Inches(5.0), width=Inches(10), height=Inches(1.7), title_color=TEAL)


# ════════════════════════════════════════════════════════════
# SLIDE 8 — P4AIR BASELINE ARCHITECTURE
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 8, TOTAL_SLIDES)
add_title_box(sl, "Baseline: P4air Architecture", "Three-module CCA-aware pipeline (Turkovic & Kuipers, 2021)")

add_info_card(sl, "Module 1: Fingerprinting", [
    "• SYN-based RTT estimation",
    "• Per-RTT packet count tracking",
    "• Queue-depth growth streak analysis",
    "• BBR bandwidth probing detection",
    "",
    "Flow lifecycle:",
    "Ant → Mice → Delay → Loss-Delay → Pure Loss",
    "(+ Model-Based branch for BBR)",
], left=Inches(0.5), top=Inches(1.6), width=Inches(3.8), height=Inches(3.5))

add_info_card(sl, "Module 2: Reallocation", [
    "• 8 hardware queues (Q0–Q7)",
    "• Q0, Q1 reserved for short flows",
    "• Q2–Q7 dynamically allocated",
    "• Proportional to active flows per group",
    "",
    "Boundaries l₁, l₂, l₃ updated",
    "each RTT interval",
], left=Inches(4.6), top=Inches(1.6), width=Inches(3.8), height=Inches(3.5))

add_info_card(sl, "Module 3: Apply Actions", [
    "When flow exceeds BDP fair share:",
    "",
    "• Loss-based → DROP packet",
    "• Delay-based → RECIRCULATE (delay)",
    "• Model-based → HALVE TCP window",
    "",
    "CCA-specific penalties = precise",
    "but complex to implement",
], left=Inches(8.7), top=Inches(1.6), width=Inches(3.8), height=Inches(3.5))

# Bottom note
add_info_card(sl, "⚠  P4air in BMv2: 13+ registers per flow, RTT-dependent, recirculation-heavy", [
], left=Inches(0.5), top=Inches(5.5), width=Inches(12), height=Inches(0.8), title_color=RED_C)


# ════════════════════════════════════════════════════════════
# SLIDE 9 — P4AIR LIMITATIONS IN BMV2
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 9, TOTAL_SLIDES)
add_title_box(sl, "P4air: Software Emulation Challenges", "Why complexity hurts in BMv2")

add_info_card(sl, "1. CPU Timing Jitter", [
    "P4air needs μs-accurate timestamps",
    "for RTT estimation.",
    "",
    "BMv2 on VM CPU introduces random",
    "delays → flow misclassification.",
    "",
    "A CUBIC flow may be wrongly labelled",
    "as delay-based (or vice versa).",
], left=Inches(0.5), top=Inches(1.5), width=Inches(3.8), height=Inches(3.5))

add_info_card(sl, "2. Recirculation Bottleneck", [
    "Delay action uses recirculate() to",
    "push packets back through pipeline.",
    "",
    "On ASIC hardware: near-instant.",
    "On BMv2 software: burns CPU cycles.",
    "",
    "Throughput drops from 10 Mbps",
    "to as low as 5.52 Mbps.",
], left=Inches(4.6), top=Inches(1.5), width=Inches(3.8), height=Inches(3.5))

add_info_card(sl, "3. Scale Dependency", [
    "With only 4 flows across 8 queues,",
    "hash-based separation works by luck.",
    "",
    "Each flow lands in its own queue →",
    "artificial fairness.",
    "",
    "P4air's real advantage only shows",
    "at 16+ flows (hash collisions).",
], left=Inches(8.7), top=Inches(1.5), width=Inches(3.8), height=Inches(3.5))

# Bottom takeaway
add_info_card(sl, "Takeaway:  Simpler dataplane logic can be MORE robust under software emulation constraints", [
], left=Inches(0.5), top=Inches(5.5), width=Inches(12), height=Inches(0.8), title_color=GREEN)


# ════════════════════════════════════════════════════════════
# SLIDE 10 — KBCS CORE IDEA
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(sl)
# Section divider
t = sl.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11), Inches(2.5))
tf = t.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Our Proposal"; p.font.size = Pt(22); p.font.color.rgb = RGBColor(150, 190, 230); p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Karma-Based Credit Scheduler"; p2.font.size = Pt(40); p2.font.bold = True; p2.font.color.rgb = WHITE; p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(12)
p3 = tf.add_paragraph(); p3.text = "\"Know if a flow is aggressive — not which CCA it runs\""; p3.font.size = Pt(20); p3.font.italic = True; p3.font.color.rgb = RGBColor(150, 220, 200); p3.alignment = PP_ALIGN.CENTER; p3.space_before = Pt(20)
# bar
bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.3), Inches(5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()
add_slide_number(sl, 10, TOTAL_SLIDES)


# ════════════════════════════════════════════════════════════
# SLIDE 11 — KBCS DESIGN PHILOSOPHY
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 11, TOTAL_SLIDES)
add_title_box(sl, "KBCS: Design Philosophy", "Three guiding principles")

add_info_card(sl, "1. Behavior Over Identity", [
    "Don't classify the CCA algorithm.",
    "Just measure: is this flow sending",
    "more than its fair share right now?",
    "",
    "→  Simpler logic, fewer registers",
    "→  No RTT estimation needed",
    "→  No recirculation required",
], left=Inches(0.5), top=Inches(1.6), width=Inches(3.8), height=Inches(3.5))

add_info_card(sl, "2. Fast Penalty, Slow Recovery", [
    "Aggressive flows lose karma quickly.",
    "Well-behaved flows regain trust slowly.",
    "",
    "Penalty:    –5 per aggressive packet",
    "Reward:    +1 per compliant packet",
    "",
    "→  Asymmetry prevents oscillation",
    "→  Stabilizes queue behavior",
], left=Inches(4.6), top=Inches(1.6), width=Inches(3.8), height=Inches(3.5))

add_info_card(sl, "3. Queue-Based Enforcement", [
    "Map karma to priority queues:",
    "",
    "  GREEN  (K>80)  → High priority",
    "  YELLOW (40<K≤80) → Medium",
    "  RED     (K≤40)  → Low / Drop",
    "",
    "→  Protects well-behaved flows",
    "→  Forces aggressors to slow down",
], left=Inches(8.7), top=Inches(1.6), width=Inches(3.8), height=Inches(3.5))

# Bottom contrast
tb = sl.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(0.8))
tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "P4air asks: \"What CCA is this flow running?\"     →     KBCS asks: \"Is this flow being greedy?\""
p.font.size = Pt(18); p.font.color.rgb = TEAL; p.font.bold = True; p.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════
# SLIDE 12 — KBCS ALGORITHM STEP-BY-STEP
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 12, TOTAL_SLIDES)
add_title_box(sl, "The Karma Algorithm", "What happens when a packet arrives at the switch")

steps = [
    ("Step 1: Flow Identification", "Hash the 5-tuple (src/dst IP, ports, protocol) using CRC16\n→ flow_id indexes into 65,536-entry register array"),
    ("Step 2: Decayed Byte Estimation", "flow_bytes = (flow_bytes >> 1) + packet_length\n→ EWMA approximation using only bit-shift (no floating point, no timers)"),
    ("Step 3: Karma Update", "If flow_bytes > threshold → Karma -= 5  (penalty)\nIf flow_bytes ≤ threshold → Karma += 1  (reward)\nClamp to range [0, 100]"),
    ("Step 4: Color Classification", "K > 80 → GREEN  |  40 < K ≤ 80 → YELLOW  |  K ≤ 40 → RED"),
    ("Step 5: Queue Mapping + AQM", "GREEN→ high-priority queue  |  YELLOW→ medium  |  RED→ low / drop"),
]
for i, (title, desc) in enumerate(steps):
    y = Inches(1.5) + Inches(1.02) * i
    # step box
    box = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), y, Inches(12.3), Inches(0.92))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(245, 248, 252) if i % 2 == 0 else WHITE
    box.line.color.rgb = RGBColor(200, 210, 225); box.line.width = Pt(0.5)
    tf = box.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = BLUE
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13); p2.font.color.rgb = GRAY


# ════════════════════════════════════════════════════════════
# SLIDE 13 — KARMA STATE MACHINE
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 13, TOTAL_SLIDES)
add_title_box(sl, "Karma State Transitions", "Fast penalty, slow recovery — visualized")

# Green state
g = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.5), Inches(2.5), Inches(2.8), Inches(2.0))
g.fill.solid(); g.fill.fore_color.rgb = RGBColor(200, 240, 210)
g.line.color.rgb = GREEN; g.line.width = Pt(2)
tf = g.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "GREEN"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = GREEN; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "K > 80"; p2.font.size = Pt(14); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph(); p3.text = "High Priority"; p3.font.size = Pt(12); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER

# Yellow state
y_shape = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.3), Inches(2.5), Inches(2.8), Inches(2.0))
y_shape.fill.solid(); y_shape.fill.fore_color.rgb = RGBColor(255, 240, 200)
y_shape.line.color.rgb = ORANGE; y_shape.line.width = Pt(2)
tf = y_shape.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "YELLOW"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = ORANGE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "40 < K ≤ 80"; p2.font.size = Pt(14); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph(); p3.text = "Medium Priority"; p3.font.size = Pt(12); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER

# Red state
r = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.1), Inches(2.5), Inches(2.8), Inches(2.0))
r.fill.solid(); r.fill.fore_color.rgb = RGBColor(255, 210, 210)
r.line.color.rgb = RED_C; r.line.width = Pt(2)
tf = r.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "RED"; p.font.size = Pt(20); p.font.bold = True; p.font.color.rgb = RED_C; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "K ≤ 40"; p2.font.size = Pt(14); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER
p3 = tf.add_paragraph(); p3.text = "Low / Drop"; p3.font.size = Pt(12); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER

# Arrows annotations (text-based since we can't draw curved arrows easily)
# Penalty arrows (top)
arr_t = sl.shapes.add_textbox(Inches(3.2), Inches(2.0), Inches(3.0), Inches(0.5))
p = arr_t.text_frame.paragraphs[0]; p.text = "──── Aggressive (K -= 5) ────▶"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = RED_C; p.alignment = PP_ALIGN.CENTER
arr_t2 = sl.shapes.add_textbox(Inches(7.0), Inches(2.0), Inches(3.0), Inches(0.5))
p = arr_t2.text_frame.paragraphs[0]; p.text = "──── Aggressive (K -= 5) ────▶"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = RED_C; p.alignment = PP_ALIGN.CENTER

# Recovery arrows (bottom)
arr_b = sl.shapes.add_textbox(Inches(3.2), Inches(4.7), Inches(3.0), Inches(0.5))
p = arr_b.text_frame.paragraphs[0]; p.text = "◀──── Compliant (K += 1) ────"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = GREEN; p.alignment = PP_ALIGN.CENTER
arr_b2 = sl.shapes.add_textbox(Inches(7.0), Inches(4.7), Inches(3.0), Inches(0.5))
p = arr_b2.text_frame.paragraphs[0]; p.text = "◀──── Compliant (K += 1) ────"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = GREEN; p.alignment = PP_ALIGN.CENTER

# New flow annotation
nf = sl.shapes.add_textbox(Inches(1.7), Inches(1.5), Inches(2.5), Inches(0.5))
p = nf.text_frame.paragraphs[0]; p.text = "↓ New Flow (K = 100)"; p.font.size = Pt(14); p.font.italic = True; p.font.color.rgb = GRAY; p.alignment = PP_ALIGN.CENTER

# Bottom note
note = sl.shapes.add_textbox(Inches(1), Inches(5.5), Inches(11), Inches(1.0))
ntf = note.text_frame; ntf.word_wrap = True
p = ntf.paragraphs[0]; p.text = "Recovery math:  RED → YELLOW requires 40+ consecutive compliant packets"; p.font.size = Pt(15); p.font.color.rgb = BLUE; p.font.bold = True; p.alignment = PP_ALIGN.CENTER
p2 = ntf.add_paragraph(); p2.text = "YELLOW → GREEN requires another 20+ compliant packets.   This asymmetry prevents oscillation."; p2.font.size = Pt(14); p2.font.color.rgb = GRAY; p2.alignment = PP_ALIGN.CENTER


# ════════════════════════════════════════════════════════════
# SLIDE 14 — KBCS PROCESSING PIPELINE
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 14, TOTAL_SLIDES)
add_title_box(sl, "KBCS Processing Pipeline", "End-to-end packet flow through the programmable switch")

# Pipeline stages as cards
stages = [
    ("Senders", "h1: CUBIC\nh2: BBR\n(100 Mbps each)", RGBColor(200, 220, 240), BLUE),
    ("Parser &\n5-Tuple Hash", "Extract headers\nCRC16 → flow_id\n(65K register array)", RGBColor(220, 240, 230), TEAL),
    ("Decayed Byte\nTracking", "flow_bytes =\n(old >> 1) + pkt_len\n(EWMA, no timer)", RGBColor(220, 240, 230), TEAL),
    ("Karma\nUpdate", "Aggressive? K -= 5\nCompliant? K += 1\nClamp [0,100]", RGBColor(220, 240, 230), TEAL),
    ("Flow Color\n& Queue Map", "GREEN → Q_high\nYELLOW → Q_med\nRED → Q_low/Drop", RGBColor(220, 240, 230), TEAL),
    ("Bottleneck\n+ Receiver", "10 Mbps link\n5ms delay\n→ h3 receiver", RGBColor(240, 220, 220), RED_C),
]
for i, (title, body, bg_color, title_clr) in enumerate(stages):
    x = Inches(0.3) + Inches(2.1) * i
    card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.8), Inches(1.95), Inches(3.0))
    card.fill.solid(); card.fill.fore_color.rgb = bg_color
    card.line.color.rgb = title_clr; card.line.width = Pt(1.5)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1); tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = title_clr; p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph(); p2.text = ""; p2.font.size = Pt(4)
    p3 = tf.add_paragraph(); p3.text = body; p3.font.size = Pt(11); p3.font.color.rgb = GRAY; p3.alignment = PP_ALIGN.CENTER
    # Arrow between cards
    if i < len(stages) - 1:
        arr = sl.shapes.add_textbox(x + Inches(1.95), Inches(3.0), Inches(0.2), Inches(0.4))
        ap = arr.text_frame.paragraphs[0]; ap.text = "▶"; ap.font.size = Pt(16); ap.font.color.rgb = TEAL; ap.alignment = PP_ALIGN.CENTER

# Key advantage box
add_info_card(sl, "Key Advantage Over P4air", [
    "• Only 2 registers per flow  (vs. 13+ for P4air)",
    "• Zero recirculation  (no CPU overhead in BMv2)",
    "• No timer dependency  (decay happens on packet arrival)",
    "• O(1) per-packet processing  (constant-time lookup and update)",
], left=Inches(0.5), top=Inches(5.1), width=Inches(12), height=Inches(1.5), title_color=GREEN)


# ════════════════════════════════════════════════════════════
# SLIDE 15 — HOW SENDERS RESPOND
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 15, TOTAL_SLIDES)
add_title_box(sl, "Feedback Loop: How Senders React", "KBCS influences senders through standard TCP congestion signals")

add_info_card(sl, "CUBIC (Loss-Based — The Aggressor)", [
    "KBCS drops packets when Karma ≤ 40 (RED)",
    "",
    "CUBIC detects loss via duplicate ACKs",
    "→ Performs multiplicative decrease",
    "→ Cuts CWND by ~30%",
    "",
    "This IS the intended behavior.",
    "KBCS forces CUBIC to see the congestion",
    "it's causing.",
], left=Inches(0.5), top=Inches(1.6), width=Inches(3.8), height=Inches(4.2))

add_info_card(sl, "BBR (Model-Based — The Prober)", [
    "KBCS places low-karma BBR flows in",
    "lower-priority queues → increased RTT",
    "",
    "BBR observes increased RTT",
    "→ Adjusts pacing rate downward",
    "",
    "Queue-based delay is MORE effective",
    "for BBR than packet drops, since BBR",
    "deliberately ignores isolated losses.",
], left=Inches(4.6), top=Inches(1.6), width=Inches(3.8), height=Inches(4.2))

add_info_card(sl, "Vegas (Delay-Based — The Victim)", [
    "Vegas naturally backs off when RTT rises.",
    "",
    "KBCS PROTECTS Vegas flows:",
    "→ They stay GREEN (high karma)",
    "→ High-priority queue",
    "→ Minimal queuing delay",
    "",
    "Result: Vegas gets fair bandwidth share",
    "instead of being starved.",
], left=Inches(8.7), top=Inches(1.6), width=Inches(3.8), height=Inches(4.2))


# ════════════════════════════════════════════════════════════
# SLIDE 16 — P4AIR vs KBCS COMPARISON
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 16, TOTAL_SLIDES)
add_title_box(sl, "P4air vs KBCS: Side-by-Side", "Precision vs. simplicity tradeoff")

rows = [
    ["Aspect", "P4air", "KBCS"],
    ["Decision Basis", "Full CCA fingerprinting", "Behavior-based karma score"],
    ["State Per Flow", "13+ registers", "2 registers (bytes, karma)"],
    ["Classification", "4 CCA groups + lifecycle", "3 colors (GREEN/YELLOW/RED)"],
    ["Penalty Type", "CCA-specific (drop/delay/halve)", "Uniform karma degradation"],
    ["Recovery", "Implicit (reclassification)", "Explicit (+1 per packet)"],
    ["Recirculation", "Required", "Not needed"],
    ["Timer Dependency", "Critical (RTT estimation)", "None (packet-driven decay)"],
    ["BMv2 Robustness", "High variance (σ = ±1.61 Mbps)", "Expected lower variance"],
]
make_table(sl, rows, Inches(0.6), Inches(1.6), [Inches(2.8), Inches(4.5), Inches(4.5)], Inches(0.53))


# ════════════════════════════════════════════════════════════
# SLIDE 17 — CURRENT IMPLEMENTATION STATUS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 17, TOTAL_SLIDES)
add_title_box(sl, "What's Built So Far", "Current repository status at mid-semester")

# Done items
done = [
    "P4air baseline: implemented & multi-run evaluated (30-run + 16-flow)",
    "KBCS design documents & SRS: comprehensive specifications complete",
    "KBCS P4 prototype: parser, hash, byte decay, karma update, color mapping",
    "Plain forwarding baseline: separate P4 program for direct comparison",
    "Mininet topology: h1(CUBIC) + h2(BBR) → s1 → h3, 10 Mbps bottleneck",
    "Experiment automation: run_experiment.py driver with iperf3 integration",
]
for i, item in enumerate(done):
    y = Inches(1.5) + Inches(0.45) * i
    mark = sl.shapes.add_textbox(Inches(0.5), y, Inches(0.5), Inches(0.4))
    mp = mark.text_frame.paragraphs[0]; mp.text = "✓"; mp.font.size = Pt(18); mp.font.bold = True; mp.font.color.rgb = GREEN
    tb = sl.shapes.add_textbox(Inches(1.0), y, Inches(6.0), Inches(0.4))
    tp = tb.text_frame.paragraphs[0]; tp.text = item; tp.font.size = Pt(14); tp.font.color.rgb = GRAY

# Pending items
pending = [
    "Priority-queue integration in main experiment flow",
    "Queue-depth-aware KBCS decision logic",
    "Comprehensive KBCS vs baseline result archive",
    "Final parameter tuning & analysis",
]
for i, item in enumerate(pending):
    y = Inches(4.3) + Inches(0.45) * i
    mark = sl.shapes.add_textbox(Inches(7.5), y, Inches(0.5), Inches(0.4))
    mp = mark.text_frame.paragraphs[0]; mp.text = "○"; mp.font.size = Pt(18); mp.font.bold = True; mp.font.color.rgb = ORANGE
    tb = sl.shapes.add_textbox(Inches(8.0), y, Inches(5.0), Inches(0.4))
    tp = tb.text_frame.paragraphs[0]; tp.text = item; tp.font.size = Pt(14); tp.font.color.rgb = GRAY

# Labels
lbl1 = sl.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(3), Inches(0.4))
p = lbl1.text_frame.paragraphs[0]; p.text = "COMPLETED"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = GREEN
lbl2 = sl.shapes.add_textbox(Inches(7.5), Inches(3.9), Inches(3), Inches(0.4))
p = lbl2.text_frame.paragraphs[0]; p.text = "IN PROGRESS / PENDING"; p.font.size = Pt(16); p.font.bold = True; p.font.color.rgb = ORANGE


# ════════════════════════════════════════════════════════════
# SLIDE 18 — PROTOTYPE PARAMETERS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 18, TOTAL_SLIDES)
add_title_box(sl, "KBCS Prototype Parameters", "Current values implemented in P4 code")

rows = [
    ["Parameter", "Value", "Purpose"],
    ["Register Size", "65,536", "Max concurrent flows (CRC16 hash space)"],
    ["Initial Karma", "100", "All flows start fully trusted"],
    ["Penalty Step", "–5", "Fast punishment for aggressive sending"],
    ["Reward Step", "+1", "Gradual trust recovery"],
    ["High Threshold (K)", "80", "GREEN boundary"],
    ["Low Threshold (K)", "40", "RED boundary"],
    ["Byte Threshold (θ)", "120,000 bytes", "Aggressiveness detection level"],
    ["Decay Formula", "bytes = (bytes >> 1) + pkt_len", "EWMA via bit-shift (50% decay)"],
]
make_table(sl, rows, Inches(0.8), Inches(1.6), [Inches(2.8), Inches(3.0), Inches(5.8)], Inches(0.55))

add_info_card(sl, "Design Gap — Honest Assessment", [
    "Current prototype ≠ full enhanced design target.",
    "",
    "• Queue-depth-aware decision logic → not yet integrated",
    "• Selective AQM (drop only critical RED) → current code",
    "   drops all RED packets immediately",
    "• Priority queue flag → not passed consistently in",
    "   experiment runner",
    "",
    "This is acknowledged and planned for Phase 2.",
], left=Inches(0.8), top=Inches(5.0), width=Inches(11.5), height=Inches(1.6), title_color=ORANGE)


# ════════════════════════════════════════════════════════════
# SLIDE 19 — BASELINE RESULTS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 19, TOTAL_SLIDES)
add_title_box(sl, "Baseline Evaluation Results", "Measured, reproducible, already in the repository")

# 30-run table
t1 = sl.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(6), Inches(0.4))
p = t1.text_frame.paragraphs[0]; p.text = "30-Run Average (4 flows: CUBIC + BBR + Vegas + Illinois)"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = BLUE

rows1 = [
    ["Configuration", "Avg Jain's FI", "Avg Throughput", "Std Dev (Mbps)"],
    ["No AQM (FIFO)", "0.9672", "9.81 Mbps", "± 0.27"],
    ["Diff Queues (Hash)", "0.9541", "10.02 Mbps", "± 0.43"],
    ["P4air (CCA Aware)", "0.9452", "9.68 Mbps", "± 1.61"],
]
make_table(sl, rows1, Inches(0.5), Inches(1.9), [Inches(2.5), Inches(2.0), Inches(2.2), Inches(2.0)], Inches(0.45))

# 16-flow table
t2 = sl.shapes.add_textbox(Inches(0.5), Inches(4.1), Inches(6), Inches(0.4))
p = t2.text_frame.paragraphs[0]; p.text = "16-Flow Scale Test (4× each CCA)"; p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = BLUE

rows2 = [
    ["Configuration", "Avg Jain's FI", "Avg Throughput", "Std Dev (Mbps)"],
    ["No AQM (FIFO)", "0.9123", "10.23 Mbps", "± 0.03"],
    ["Diff Queues (Hash)", "0.9390", "10.35 Mbps", "± 0.02"],
    ["P4air (CCA Aware)", "0.9483", "11.02 Mbps", "± 0.02"],
]
make_table(sl, rows2, Inches(0.5), Inches(4.6), [Inches(2.5), Inches(2.0), Inches(2.2), Inches(2.0)], Inches(0.45))

# Right side insights
add_info_card(sl, "Key Observations", [
    "At 4 flows: FIFO appears fair (artifactual)",
    "  → Lucky hash effect at small scale",
    "  → P4air pays CPU cost for limited gain",
    "",
    "At 16 flows: P4air truly shines",
    "  → FIFO fairness drops (0.97 → 0.91)",
    "  → P4air leads in both FI and throughput",
    "",
    "P4air's variance is concerning:",
    "  σ = ±1.61 Mbps (recirculation cost)",
    "  → KBCS aims for σ < ±0.5 Mbps",
], left=Inches(9.4), top=Inches(1.6), width=Inches(3.5), height=Inches(4.5))


# ════════════════════════════════════════════════════════════
# SLIDE 20 — KBCS PERFORMANCE TARGETS
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 20, TOTAL_SLIDES)
add_title_box(sl, "KBCS Performance Targets", "What we need to achieve for a credible evaluation")

add_info_card(sl, "Target 1: Fairness", [
    "Jain's FI ≥ 0.95 at 16-flow scale",
    "",
    "Must match or exceed FIFO (0.91)",
    "and approach P4air (0.95)",
], left=Inches(0.5), top=Inches(1.6), width=Inches(3.8), height=Inches(2.2), title_color=GREEN)

add_info_card(sl, "Target 2: Utilization", [
    "Link utilization ≥ 90%",
    "",
    "Fairness should not collapse throughput.",
    "PFQ shows FI > 0.95 at U > 95% is possible.",
], left=Inches(4.6), top=Inches(1.6), width=Inches(3.8), height=Inches(2.2), title_color=TEAL)

add_info_card(sl, "Target 3: Stability", [
    "Throughput σ < ± 0.5 Mbps",
    "",
    "P4air has σ = ± 1.61 Mbps in BMv2.",
    "Simpler logic should be more stable.",
], left=Inches(8.7), top=Inches(1.6), width=Inches(3.8), height=Inches(2.2), title_color=BLUE)

# Evaluation methodology
add_info_card(sl, "Evaluation Methodology", [
    "• Dumbbell topology: 10 Mbps bottleneck, 5ms delay",
    "• iperf3 traffic generation with controlled CCA mixtures",
    "• Jain's Fairness Index:  J = (Σxᵢ)² / (n · Σxᵢ²)",
    "• 30+ repeated runs for statistical significance",
    "• Compare: FIFO  vs  Diff Queues  vs  P4air  vs  KBCS",
], left=Inches(0.5), top=Inches(4.2), width=Inches(12), height=Inches(2.2))


# ════════════════════════════════════════════════════════════
# SLIDE 21 — ROADMAP
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(sl); add_bottom_bar(sl); add_slide_number(sl, 21, TOTAL_SLIDES)
add_title_box(sl, "End-Semester Roadmap", "Four phases from prototype to final evaluation")

phases = [
    ("Phase 1", "Complete Queue-Aware Logic", [
        "Integrate queue-depth-aware decisions",
        "Align parameters with enhanced design",
        "Selective AQM for RED flows",
    ], GREEN),
    ("Phase 2", "Finalize Experiment Pipeline", [
        "BMv2 priority queues in all KBCS runs",
        "Archive KBCS vs baseline results",
        "Automated plotting scripts",
    ], TEAL),
    ("Phase 3", "Comparative Evaluation", [
        "KBCS vs FIFO vs Hash vs P4air",
        "30+ repeated runs per config",
        "Analysis: fairness, throughput, variance",
    ], BLUE),
    ("Phase 4", "Final Report & Analysis", [
        "Insert final KBCS result plots",
        "Sensitivity analysis (flow count, thresholds)",
        "Conclusions based on measured outcomes",
    ], RGBColor(100, 70, 130)),
]

for i, (phase, title, items, color) in enumerate(phases):
    x = Inches(0.4) + Inches(3.15) * i
    # Phase header
    hdr = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(1.6), Inches(2.95), Inches(0.6))
    hdr.fill.solid(); hdr.fill.fore_color.rgb = color
    hdr.line.fill.background()
    tf = hdr.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = f"{phase}: {title}"; p.font.size = Pt(13); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    # Items card
    card = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Inches(2.3), Inches(2.95), Inches(2.5))
    card.fill.solid(); card.fill.fore_color.rgb = RGBColor(245, 248, 252)
    card.line.color.rgb = color; card.line.width = Pt(1.5)
    tf = card.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.12); tf.margin_top = Inches(0.1)
    for j, item in enumerate(items):
        pp = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        pp.text = f"• {item}"; pp.font.size = Pt(12); pp.font.color.rgb = GRAY; pp.space_after = Pt(6)
    # Arrow
    if i < len(phases) - 1:
        arr = sl.shapes.add_textbox(x + Inches(2.95), Inches(2.8), Inches(0.25), Inches(0.4))
        ap = arr.text_frame.paragraphs[0]; ap.text = "▶"; ap.font.size = Pt(16); ap.font.color.rgb = TEAL

# Lessons learned
add_info_card(sl, "Lessons Learned So Far", [
    "1. Evaluate at multiple traffic scales — small-scale results can be misleading",
    "2. Complexity in BMv2 causes high variance — simpler can be better",
    "3. A clean baseline is essential before claiming any improvement",
    "4. Current prototype validates feasibility — full evaluation needs all phases",
], left=Inches(0.4), top=Inches(5.1), width=Inches(12.5), height=Inches(1.5))


# ════════════════════════════════════════════════════════════
# SLIDE 22 — THANK YOU / Q&A
# ════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(sl)

bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4), Inches(2.9), Inches(5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()

t = sl.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11), Inches(1.5))
tf = t.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; p.text = "Thank You"; p.font.size = Pt(48); p.font.bold = True; p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph(); p2.text = "Questions & Discussion"; p2.font.size = Pt(24); p2.font.color.rgb = RGBColor(150, 190, 230); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(10)

# Summary box
summary = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2), Inches(3.5), Inches(9), Inches(2.8))
summary.fill.solid(); summary.fill.fore_color.rgb = RGBColor(25, 70, 120)
summary.line.color.rgb = TEAL; summary.line.width = Pt(1.5)
tf = summary.text_frame; tf.word_wrap = True
tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)

items = [
    ("Key Takeaways:", Pt(18), WHITE, True),
    ("", Pt(6), WHITE, False),
    ("1.  TCP fairness at the bottleneck is a real, measurable problem", Pt(15), RGBColor(200, 220, 240), False),
    ("2.  P4 enables in-switch behavior tracking — we proved this with P4air baseline", Pt(15), RGBColor(200, 220, 240), False),
    ("3.  KBCS's \"know behavior, not CCA\" approach = simpler & potentially more robust", Pt(15), RGBColor(200, 220, 240), False),
    ("4.  Prototype working, baseline validated, end-semester evaluation planned", Pt(15), RGBColor(200, 220, 240), False),
]
for i, (txt, sz, clr, bold) in enumerate(items):
    pp = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    pp.text = txt; pp.font.size = sz; pp.font.color.rgb = clr; pp.font.bold = bold; pp.space_after = Pt(4)

# Contact info
contact = sl.shapes.add_textbox(Inches(2), Inches(6.5), Inches(9), Inches(0.6))
ctf = contact.text_frame; ctf.word_wrap = True
cp = ctf.paragraphs[0]
cp.text = "Gaurav Jaiswal  ·  Ishan Chadha  ·  Sameer Prasad   |   IIIT Allahabad   |   March 2026"
cp.font.size = Pt(13); cp.font.color.rgb = RGBColor(150, 170, 200); cp.alignment = PP_ALIGN.CENTER

# ── Save ────────────────────────────────────────────────────
output_path = r"e:\Research Methodology\Project-Implementation\KBCS_Mid_Semester_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
